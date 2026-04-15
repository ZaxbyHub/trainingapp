"""
Document Processor Module
Handles extraction and chunking of PDF, DOCX, and PPTX files.
"""

import os
import re
import logging
from typing import List, Tuple, Optional
from dataclasses import dataclass
from pathlib import Path

logger = logging.getLogger(__name__)

ABBREVIATIONS = frozenset({
    'dr', 'mr', 'mrs', 'ms', 'prof', 'jr', 'sr', 'st', 'ave', 'blvd',
    'dept', 'rev', 'vol', 'fig', 'ed', 'eds', 'repr', 'trans', 'pt',
    'ch', 'sec', 'app', 'ex', 'cf', 'eg', 'ie', 'etc', 'approx',
    'esp', 'viz', 'al', 'vs', 'inc', 'corp', 'ltd', 'govt', 'est',
    'acct', 'tel', 'ref',
})


@dataclass
class DocumentChunk:
    """Represents a chunk of text from a document."""

    text: str
    source: str
    page: Optional[int] = None
    chunk_index: int = 0


class DocumentProcessor:
    """Processes various document formats and extracts text."""

    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".md"}

    def __init__(self, chunk_size: int = 256, chunk_overlap: int = 50):
        if chunk_size <= 0:
            raise ValueError(f"chunk_size must be positive, got {chunk_size}")
        if chunk_overlap < 0:
            raise ValueError(f"chunk_overlap must be non-negative, got {chunk_overlap}")
        if chunk_overlap >= chunk_size:
            raise ValueError(
                f"chunk_overlap ({chunk_overlap}) must be less than chunk_size ({chunk_size})"
            )
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def extract_pdf(self, filepath: str) -> Tuple[str, List[Tuple[int, str]]]:
        """Extract text from PDF with page information."""
        try:
            import pdfplumber

            pages_text = []
            full_text = []

            with pdfplumber.open(filepath) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    text = page.extract_text() or ""
                    if text.strip():
                        pages_text.append((i, text))
                        full_text.append(text)

            return "\n\n".join(full_text), pages_text
        except ImportError:
            logger.warning(
                "pdfplumber not installed. Falling back to pypdf for PDF extraction. "
                "Consider installing pdfplumber for better extraction quality: pip install pdfplumber"
            )
            from pypdf import PdfReader

            reader = PdfReader(filepath)
            pages_text = []
            full_text = []

            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    pages_text.append((i, text))
                    full_text.append(text)

            return "\n\n".join(full_text), pages_text

    def _extract_table_rows(self, table) -> List[str]:
        """Extract text from table rows."""
        rows = []
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                rows.append(" | ".join(row_text))
        return rows

    def extract_docx(self, filepath: str) -> str:
        """Extract text from DOCX file."""
        from docx import Document

        doc = Document(filepath)
        paragraphs = []

        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        for table in doc.tables:
            paragraphs.extend(self._extract_table_rows(table))

        return "\n\n".join(paragraphs)

    def extract_pptx(self, filepath: str) -> str:
        """Extract text from PPTX file."""
        from pptx import Presentation

        prs = Presentation(filepath)
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = [f"[Slide {slide_num}]"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text)

                if hasattr(shape, "has_table") and shape.has_table:
                    slide_content.extend(self._extract_table_rows(shape.table))

            if len(slide_content) > 1:
                slides_text.append("\n".join(slide_content))

        return "\n\n".join(slides_text)

    def extract_text_file(self, filepath: str) -> str:
        """Extract text from plain text or markdown files."""
        encodings = ["utf-8", "utf-16", "latin-1", "cp1252"]

        for encoding in encodings:
            try:
                with open(filepath, "r", encoding=encoding) as f:
                    return f.read()
            except (UnicodeDecodeError, UnicodeError):
                continue

        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    def extract_document(self, filepath: str) -> Tuple[str, List[Tuple[int, str]]]:
        """Extract text from any supported document type.
        Returns (full_text, pages) where pages is [(page_num, page_text), ...].
        For non-PDF formats, pages is empty.
        """
        filepath = str(filepath)
        ext = Path(filepath).suffix.lower()

        if ext == ".pdf":
            return self.extract_pdf(filepath)
        elif ext in {".docx", ".doc"}:
            return self.extract_docx(filepath), []
        elif ext in {".pptx", ".ppt"}:
            return self.extract_pptx(filepath), []
        elif ext in {".txt", ".md"}:
            return self.extract_text_file(filepath), []
        else:
            raise ValueError(f"Unsupported file format: {ext}")

    def clean_text(self, text: str) -> str:
        """Clean and normalize text."""
        text = re.sub(r"\s+", " ", text)
        text = re.sub(r"\n\s*\n", "\n\n", text)
        text = text.strip()
        return text

    def _split_sentences(self, paragraph: str) -> List[str]:
        """Split paragraph into sentences, respecting common abbreviations."""
        protected = paragraph
        for abbr in ABBREVIATIONS:
            protected = re.sub(
                rf'\b{abbr}\.',
                f'{abbr}\x00',
                protected,
                flags=re.IGNORECASE,
            )
        def _protect_initial(m):
            return m.group(1) + '\x00'
        protected = re.sub(r'\b([A-Z])\.', _protect_initial, protected)
        sentences = re.split(r'(?<=[.!?])\s+', protected)
        return [s.replace('\x00', '.').strip() for s in sentences if s.strip()]

    def _calculate_overlap(
        self, sentences: List[str], overlap_size: int
    ) -> Tuple[List[str], int]:
        """Calculate sentences to keep for overlap."""
        overlap_sentences = []
        overlap_word_count = 0
        for s in reversed(sentences):
            s_word_count = len(s.split())
            if overlap_word_count + s_word_count <= overlap_size:
                overlap_sentences.insert(0, s)
                overlap_word_count += s_word_count
            else:
                break
        return overlap_sentences, overlap_word_count

    def chunk_text(self, text: str, source: str, pages: Optional[List[Tuple[int, str]]] = None) -> List[DocumentChunk]:
        """Split text into overlapping chunks respecting paragraph and sentence boundaries."""
        text = self.clean_text(text)
        # Build page mapping from PDF pages
        para_page_map: dict = {}
        if pages:
            for page_num, page_text in pages:
                cleaned = re.sub(r"\s+", " ", page_text.strip())
                for segment in cleaned.split("\n\n"):
                    seg = segment.strip()
                    if seg:
                        para_page_map[seg[:80]] = page_num

        def _find_page(chunk_text: str) -> Optional[int]:
            """Find the page number for a chunk text using prefix matching."""
            if not para_page_map:
                return None
            # Try longest prefix match first
            for length in range(len(chunk_text), 0, -1):
                prefix = chunk_text[:length].strip()
                if prefix in para_page_map:
                    return para_page_map[prefix]
            return None

        paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]

        chunks = []
        chunk_index = 0
        current_chunk_sentences = []
        current_chunk_word_count = 0

        for paragraph in paragraphs:
            sentences = self._split_sentences(paragraph)

            for sentence in sentences:
                sentence = sentence.strip()
                if not sentence:
                    continue

                sentence_word_count = len(sentence.split())

                # If sentence alone exceeds chunk_size and we're starting fresh
                if (
                    sentence_word_count > self.chunk_size
                    and not current_chunk_sentences
                ):
                    # Split sentence into word chunks
                    words = sentence.split()
                    while words:
                        chunk_words = words[: self.chunk_size]
                        chunk_text = " ".join(chunk_words)
                        chunks.append(
                            DocumentChunk(
                                text=chunk_text, source=source, chunk_index=chunk_index,
                                page=_find_page(chunk_text),
                            )
                        )
                        chunk_index += 1
                        words = words[self.chunk_size :]
                        if words:
                            # Overlap handling for split sentence
                            overlap_words = (
                                chunk_words[-self.chunk_overlap // 2 :]
                                if self.chunk_overlap > 0
                                else []
                            )
                            words = overlap_words + words
                    continue  # Move to next sentence

                # If adding this sentence would exceed chunk_size, finalize current chunk
                if current_chunk_sentences and (
                    current_chunk_word_count + sentence_word_count > self.chunk_size
                ):
                    chunk_text = " ".join(current_chunk_sentences)
                    chunks.append(
                        DocumentChunk(
                            text=chunk_text, source=source, chunk_index=chunk_index,
                            page=_find_page(chunk_text),
                        )
                    )
                    chunk_index += 1

                    # Handle overlap using helper method
                    overlap_sentences, overlap_word_count = self._calculate_overlap(
                        current_chunk_sentences, self.chunk_overlap
                    )
                    current_chunk_sentences = overlap_sentences
                    current_chunk_word_count = overlap_word_count

                current_chunk_sentences.append(sentence)
                current_chunk_word_count += sentence_word_count

        # Don't forget the last chunk
        if current_chunk_sentences:
            chunk_text = " ".join(current_chunk_sentences)
            chunks.append(
                DocumentChunk(text=chunk_text, source=source, chunk_index=chunk_index,
                               page=_find_page(chunk_text))
            )

        return chunks

    def process_file(
        self, filepath: str, source_name: Optional[str] = None
    ) -> List[DocumentChunk]:
        """Process a single file and return chunks."""
        filepath = str(filepath)
        # Use provided source_name or fall back to filename from path
        filename = source_name if source_name else Path(filepath).name

        try:
            text, pages = self.extract_document(filepath)
            if not text.strip():
                logger.warning("Empty content: %s", filename)
                return []

            chunks = self.chunk_text(text, filename, pages=pages)
            logger.info("Processed: %s (%d chunks)", filename, len(chunks))
            return chunks
        except ValueError as e:
            # Known error (unsupported format, etc.)
            logger.error("Failed to process %s: %s", filename, e)
            return []
        except Exception as e:
            # Unexpected error - log full exception details
            logger.exception("Unexpected error processing %s", filename)
            return []

    def process_directory(self, directory: str) -> List[DocumentChunk]:
        """Process all supported documents in a directory recursively."""
        all_chunks = []
        directory = Path(directory)

        # Get max file size from environment variable (default 100MB)
        try:
            max_file_size_mb = int(os.environ.get("RAG_MAX_FILE_SIZE", "100"))
            if max_file_size_mb <= 0:
                max_file_size_mb = 100
        except ValueError:
            max_file_size_mb = 100
        max_file_size_bytes = max_file_size_mb * 1024 * 1024

        skipped_files = []

        for root, _, files in os.walk(directory):
            for file in files:
                filepath = Path(root) / file
                ext = filepath.suffix.lower()

                if ext in self.SUPPORTED_EXTENSIONS:
                    # Check file size before processing
                    file_size = filepath.stat().st_size
                    if file_size > max_file_size_bytes:
                        size_mb = file_size / (1024 * 1024)
                        skipped_files.append((filepath.name, size_mb))
                        logger.info("Skipping %s: %.1fMB > %.0fMB limit", filepath.name, size_mb, max_file_size_mb)
                        continue

                    chunks = self.process_file(str(filepath))
                    all_chunks.extend(chunks)

        if skipped_files:
            logger.info("Skipped %d file(s) exceeding %.0fMB limit", len(skipped_files), max_file_size_mb)
        logger.info("Total: %d chunks from %s", len(all_chunks), directory)
        return all_chunks


if __name__ == "__main__":
    processor = DocumentProcessor(chunk_size=512, chunk_overlap=50)

    import sys

    if len(sys.argv) > 1:
        path = sys.argv[1]
        if os.path.isdir(path):
            chunks = processor.process_directory(path)
        else:
            chunks = processor.process_file(path)

    logger.info("Extracted %d chunks", len(chunks))
    if chunks:
        logger.info("First chunk preview:\n%s...", chunks[0].text[:200])
